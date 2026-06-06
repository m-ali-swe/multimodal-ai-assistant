import asyncio
from fastapi import FastAPI, HTTPException, Query, Request,UploadFile,File,Form,Depends,Body,Cookie,Request
from fastapi.responses import RedirectResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import JSONResponse

from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.messages import AIMessage, HumanMessage,SystemMessage,RemoveMessage
from langchain_core.runnables import RunnableConfig

from langgraph.graph import START,END,StateGraph,MessagesState
from langgraph.graph.state import CompiledStateGraph
from langgraph.prebuilt import ToolNode,tools_condition
from langgraph.checkpoint.postgres.aio import AsyncPostgresSaver
from psycopg_pool import ConnectionPool,AsyncConnectionPool

import os,json,io,uuid
from dotenv import find_dotenv, load_dotenv
from typing_extensions import Literal,List,Annotated,Optional
from contextlib import asynccontextmanager
from PyPDF2 import PdfReader
from docx import Document
# import pandas as pd
from pptx import Presentation
from src.chatbot_backend.db import get_db,Session,Chats,User #type: ignore
from src.chatbot_backend.service import get_user,signup_user #type: ignore

_:bool=load_dotenv(find_dotenv())

checkpointer: AsyncPostgresSaver | None = None
graph: CompiledStateGraph | None = None
pool: ConnectionPool | None = None
initialized = False


gemini_api_key=os.environ.get("GEMINI_API_KEY")
DB_URI=os.environ.get("DB_URI")

class MessagesState(MessagesState):
  summary:str

async def startup():
    global checkpointer, pool, graph,initialized
    
    if initialized:
        print("Startup already ran, skipping re-init.")
        return

    print("Setting up database connection...")
    connection_kwargs={"autocommit":True,"prepare_threshold":0}
    pool=AsyncConnectionPool(conninfo=DB_URI,max_size=20,kwargs=connection_kwargs)
    await pool.open()
    
    checkpointer=AsyncPostgresSaver(pool)
    await checkpointer.setup()
    print("Database connection setup complete.")
    
    # graph=workflow.compile()
    print("Application setup complete.")
    
    # initialized = True

async def get_graph():
    global checkpointer, graph
    def search_web(query: str) -> str:
        """Use this function to search the web for answers."""
        print(f"\n--- Model is calling the search_web tool with query: '{query}' ---\n")
        return f"The search for '{query}' returned the following result: The global software market is projected to reach $687.9 billion in 2025."

    model=ChatGoogleGenerativeAI(model='gemini-2.0-flash',api_key=gemini_api_key).bind_tools([search_web])
    
    def call_model(state:MessagesState,config:RunnableConfig)->MessagesState:
        summary=state.get("summary","")
        if summary:
            sys_msg=f"Summary of conversation earlier: {summary}"
            messages=[SystemMessage(content=sys_msg)]+state["messages"]
        else:
            messages=state["messages"]
        response=model.invoke(messages)
        print(f"Response : {response.content}")
        return {"messages":[response]}

    def summarize_conversation(state:MessagesState,config:RunnableConfig):
        summary=state.get("summary","")
        if summary:
            summary_message = (
                f"This is summary of the conversation to date: {summary}\n\n"
                "Extend the summary by taking into account the new messages above:"
            )
        else:
            summary_message = "Create a summary of the conversation above:"
       
        messages=state["messages"]+[HumanMessage(content=summary_message)]
        response=model.invoke(messages)
        delete_messages=[RemoveMessage(id=m.id) for m in state["messages"][:-2]]
        return {"summary":response.content,"messages":delete_messages}

    def should_continue(state: MessagesState)->Literal["summarize_conversation","tools","__end__"]:
        """Return the next node to execute."""
        if tools_condition(state)=="tools":
            return "tools"
        
        messages = state["messages"]
        
        if len(messages) > 20:
            return "summarize_conversation"
        return "__end__"

    workflow: StateGraph = StateGraph(MessagesState)
    workflow.add_node("conversation", call_model)
    workflow.add_node("summarize_conversation",summarize_conversation)
    workflow.add_node("tools",ToolNode([search_web]))

    workflow.add_edge(START, "conversation")
    workflow.add_conditional_edges("conversation", should_continue,{"summarize_conversation": "summarize_conversation",
                "tools":"tools", "__end__": END})
    workflow.add_edge("summarize_conversation", END)
    workflow.add_edge("tools","conversation")

    checkpointer=AsyncPostgresSaver(pool)
    await checkpointer.setup()
    print("Database connection setup complete.")

    
    graph=workflow.compile(checkpointer)
    return graph,checkpointer

app: FastAPI = FastAPI(
title="Gemini Chatbot Backend",
description="A simple FastAPI backend for a conversational chatbot using the Gemini API.",
)

@app.on_event("shutdown")
async def shutdown():
    global pool
    if pool:
        await pool.close()
        print("Database pool closed.")


app.add_middleware(
    CORSMiddleware,
    # allow_origins=["https://chatbot-frontend-delta-seven.vercel.app","http://localhost:3000"],
    allow_origins=[os.environ.get("ALLOWED_ORIGINS","http://localhost:3000")],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)

@app.get("/")
async def read_root():
    await startup()
    return {
        "Hello": "World",
        "graph": str(type(graph)),  # or graph.__class__.__name__
        "checkpointer": str(type(checkpointer)),
        "pool": str(type(pool))
    }

async def process_files(files:List[UploadFile]=File(...)):
    message_parts = []
    for file in files:
        file_name=file.filename
        content_type=file.content_type
        content=await file.read()
        
        # if content_type.startswith("image/"):
            
        #     file_path=save_local_file(content,file.filename)
        #     public_url = f"http://localhost:8000/uploads/{os.path.basename(file_path)}"
        #     message_parts.append({
        #         "type":"image_url",
        #         "image_url":{"url":public_url}
        #     })
                        
        if content_type.startswith("application/pdf"):
            with open(file_name,"wb") as f:
                f.write(content)

            reader=PdfReader(file_name)
            text="Data from PDF : "
            for page in reader.pages:
                text+=page.extract_text()+"\n"
            message_parts.append({
                "type":"text",
                "text":text
            })
        
        elif content_type.startswith("application/vnd.openxmlformats-officedocument.wordprocessingml.document"):
            
            doc=Document(io.BytesIO(content))
            text="Data from Docx : " + "\n".join(p.text for p in doc.paragraphs)
            message_parts.append({
                "type":"text",
                "text":text
            })
        # elif content_type.startswith("application/vnd.ms-excel") or content_type.startswith("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"):
        #     if file_name.endswith(".xlsx"):
        #         engine="openpyxl"
        #     elif file_name.endswith(".xls"):
        #         engine="xlrd"
        #     df=pd.read_excel(io.BytesIO(content),engine=engine)
        #     text="Data from Xlsx : \n" + df.to_string()
        #     message_parts.append({
        #         "type":"text",
        #         "text":text
        #     })

        elif content_type.startswith("application/vnd.openxmlformats-officedocument.presentationml.presentation"):
            presentation=Presentation(io.BytesIO(content))
            text_runs=[]
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape,"text") and shape.text.strip():
                        text_runs.append(shape.text.strip())

            text="Data from PPTx : \n" + "\n".join(text_runs)
            message_parts.append({
                "type":"text",
                "text":text
            })            
        else:
            text="Data from File : \n" + content.decode("utf-8",errors="ignore")
            message_parts.append({
                "type":"text",
                "text":text
            })
    return message_parts


@app.get("/chat/history/{thread_id}")
async def get_chat_history(request:Request,thread_id: str,db:Annotated[Session,Depends(get_db)]):
    await startup()
    user_id=request.cookies.get("user_id")
    print(f"User id in get_chat_history : {user_id} and thread : {thread_id}")
    existing_chat=db.query(Chats).filter_by(thread_id=thread_id,user_id=user_id).first()
    print(f"Existing chat : {existing_chat}")
    if not existing_chat:
        return {"messages": []}
    
    state = await checkpointer.aget({"configurable": {"thread_id": thread_id}})
    print(f"State from db : {state}")
    if not state:
        return {"messages": []}
    
    raw_messages = state["channel_values"]["messages"]
    simplified_messages = []
    for m in raw_messages:
        if isinstance(m, HumanMessage):
            content = m.content[0]

            simplified_messages.append({"role": "human", "content": content["text"]})
        
        elif isinstance(m, AIMessage):
            simplified_messages.append({"role": "ai", "content": str(m.content)})
        
        else:
            simplified_messages.append({"role": "system", "content": str(m.content)})
    print(f"Simplified messages : {simplified_messages}")
    return {"messages": simplified_messages}

@app.get("/all_chats")
async def get_all_chats(request:Request,db:Annotated[Session,Depends(get_db)],user_id:Optional[str]=Cookie(None)):
    await startup()
    user_id=request.cookies.get("user_id")
    print(f"User id in get_all_chats : {user_id}")
    if not user_id:
        raise HTTPException(status_code=401, detail="Not logged in")
    all_chats=db.query(Chats).filter_by(user_id=user_id).all()
    user=db.query(User).filter_by(id=user_id).first()
    res={"chats": all_chats, "user_id":user_id,"user_email":user.email,"total_chats":len(user.chats)}
    print(f"All chats res : {res}")
    return res

@app.post("/chat_stream")
async def chat_stream(query: str=Form(...), thread_id: str=Form(...),files:List[UploadFile]=File(None)):
    await startup()
    config = {"configurable": {"thread_id": thread_id}}
    print(f"thread_id in chat_stream : {thread_id}")
    messages = [{"type": "text", "text": query}]
    # files=[]
    if files:
        extracted_parts = await process_files(files)
        messages.extend(extracted_parts)
    graph=get_graph()
    async def event_generator():
        try:
            node_to_stream = "conversation"
            input_message = HumanMessage(content=messages)
            async for event in graph.astream_events({"messages":[input_message]}, config, version="v2"):
                if event["event"] == "on_chat_model_stream" and event["metadata"].get("langgraph_node","") == node_to_stream:
                    chunk = event["data"]["chunk"]
                    if chunk.content:
                        yield json.dumps({"type":"content","response":chunk.content}) + "\n"
        except asyncio.CancelledError:
            print("⚠️ Stream cancelled by client")
            raise
        
    
    return StreamingResponse(event_generator(),media_type="application/x-ndjson") 

@app.post("/new_chat_stream")
async def chat_stream(request:Request,db:Annotated[Session,Depends(get_db)],query: str=Form(...),files:Optional[List[UploadFile]]=File(None)):
    await startup()
    user_id=request.cookies.get("user_id")
    # user_id="1"
    print(f"User id in new_chat_stream : {user_id}")
    messages = [{"type": "text", "text": query}]
    # files=[]
    if files:
        extracted_parts = await process_files(files)
        messages.extend(extracted_parts)
    
    thread_id=str(uuid.uuid4())
    print(f"New thread id : {thread_id}")
    config = {"configurable": {"thread_id": thread_id}}
    graph=get_graph()
    async def event_generator():
        try:
            yield json.dumps({"type": "init", "thread_id": thread_id}) + "\n"
            node_to_stream = "conversation"
            input_message = HumanMessage(content=messages)
            async for event in graph.astream_events({"messages": [input_message]}, config, version="v2"):
                if event["event"] == "on_chat_model_stream" and event["metadata"].get("langgraph_node", "") == node_to_stream:
                    chunk = event["data"]["chunk"]
                    if chunk and chunk.content:
                        yield json.dumps({"type": "content", "response": chunk.content}) + "\n"
                    else:
                        yield json.dumps({"type": "keepalive"}) + "\n"
                # elif event["event"] == "on_chain_end":
                #     break  # Ensure clean exit on chain completion
        except Exception as e:
            print(f"Error in event_generator: {e}")
            yield json.dumps({"type": "error", "message": str(e)}) + "\n"
            raise

        prompt = (
            f"Return ONLY one short, descriptive title (max 5 words) for this chat. "
            f"No bullet points, no quotes, no explanations. "
            f"First message: '{query}'"
        )
        temp_model=ChatGoogleGenerativeAI(model='gemini-2.0-flash',api_key=gemini_api_key)
        print("streamed everything")        
        temp_res = await temp_model.ainvoke(prompt)
        print(f"New chat name : {temp_res.content}")
        chat_name = temp_res.content
        yield json.dumps({"type":"final","chat_name":chat_name,"thread_id":thread_id})+"\n"
        print("Yielded new chat name , going to save to db")        
        temp_chat=Chats(thread_id=thread_id,user_id=user_id,chat_name=chat_name)
        print("Adding chat to database : ", temp_chat.chat_name)
        db.add(temp_chat)
        db.commit()
        db.refresh(temp_chat)                
    # return StreamingResponse(event_generator(),media_type="application/x-ndjson",headers={"Transfer-Encoding": "chunked", "X-Accel-Buffering": "no"}) 
    return StreamingResponse(event_generator(),media_type="application/x-ndjson") 


@app.get("/rename_chat")
async def rename_chat(thread_id:str,new_name:str,db:Annotated[Session,Depends(get_db)]):
    await startup()
    chat=db.query(Chats).filter_by(thread_id=thread_id).first()
    print(f"Chat to rename : {chat} with thread_id : {thread_id}")
    chat.chat_name=new_name
    db.commit()
    db.refresh(chat)
    return {"message": "Chat renamed successfully", "chat": chat}

@app.get("/delete_chat")
async def delete_chat(thread_id:str,db:Annotated[Session,Depends(get_db)]):
    await startup()
    chat=db.query(Chats).filter_by(thread_id=thread_id).first()
    print(f"Chat to delete : {chat} with thread_id : {thread_id}")
    db.delete(chat)
    db.commit()
    return {"message": "Chat deleted successfully", "chat": chat}


@app.post("/auth/login")
def login(db:Annotated[Session,Depends(get_db)],email: str = Body(...), password: str = Body(...)):
    try:
        res:dict=get_user(db=db,email=email,password=password)
        print(f"Login response : {res}")
        if "Error" in res or "Exception" in res:
            return res
        response_data={"message": "Login successful"}
        response =JSONResponse(content=response_data)
        # response.set_cookie(key="token",value=res["access_token"],httponly=True,domain=".vercel.app",secure=True,samesite="none")
        # response.set_cookie(key="user_id",value=res["user"].id,httponly=True,domain=".vercel.app",secure=True,samesite="none")
        response.set_cookie(key="token",value=res["access_token"],httponly=True,secure=True,samesite="none")
        response.set_cookie(key="user_id",value=res["user"].id,httponly=True,secure=True,samesite="none")
        return response
    except Exception as e:
        return {"Exception":e,"type":"exception"}

@app.post("/auth/signup")
def signup(db:Annotated[Session,Depends(get_db)],email: str = Body(...), password: str = Body(...)):
    try:
        res=signup_user(db=db,email=email,password=password)
        print(f"Signup response : {res}")
        if "Error" in res or "Exception" in res:
            return res        
        response_data={"message":"account created successfully"}
        response=JSONResponse(content=response_data)
        response.set_cookie(key="token",value=res["access_token"],httponly=True,secure=True,samesite="none")
        response.set_cookie(key="user_id",value=res["new_user"].id,httponly=True,secure=True,samesite="none")
        # response.set_cookie(key="token",value=res["access_token"],httponly=True,domain=".vercel.app",secure=True,samesite="none")
        # response.set_cookie(key="user_id",value=res["new_user"].id,httponly=True,domain=".vercel.app",secure=True,samesite="none")
        return response
    except Exception as e:
        return {"error":e,"type":"execption"}


# def guest_login(db:Annotated[Session,Depends(get_db)],redirect_url: Optional[str] = Query(None)):
@app.get("/guest_login")
def guest_login(db:Annotated[Session,Depends(get_db)]):
    email=f"guest_{uuid.uuid4()}"
    password="guest"
    try:
        res=signup_user(db=db,email=email,password=password)
        print(f"Guest signup response : {res}")
        if "Error" in res or "Exception" in res:
            return res        
        # response_data={"message":"account created successfully"}
        
        response = RedirectResponse("https://chatbot-frontend-delta-seven.vercel.app/")
        # response = JSONResponse(content={"message": "account created successfully", "access_token": res["access_token"], "user_id": str(res["new_user"].id)})
        # response.set_cookie(key="token",value=res["access_token"],httponly=True,domain=".vercel.app",secure=True,samesite="none")
        # response.set_cookie(key="user_id",value=res["new_user"].id,httponly=True,domain=".vercel.app",secure=True,samesite="none")
        response.set_cookie(key="token",value=res["access_token"],httponly=True,secure=True,samesite="none")
        response.set_cookie(key="user_id",value=res["new_user"].id,httponly=True,secure=True,samesite="none")
        return response
    except Exception as e:
        # return {"error":e,"type":"execption"}
        return JSONResponse(content={"error": str(e), "type": "exception"}, status_code=500)
    
@app.get("/auth/logout")
def logout():
    response = JSONResponse(content={"message": "Logged out successfully"})
    response.delete_cookie(key="token", path="/",secure=True,samesite="none")
    response.delete_cookie(key="user_id", path="/",secure=True,samesite="none")
    # response.delete_cookie(key="token", path="/", domain=".vercel.app",secure=True,samesite="none")
    # response.delete_cookie(key="user_id", path="/", domain=".vercel.app",secure=True,samesite="none")
    print("User logged out, cookies deleted.")
    return response
