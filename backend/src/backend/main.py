from fastapi import FastAPI, UploadFile, File, Form, Depends, HTTPException, Body, Request, Cookie
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse, RedirectResponse
from langgraph.graph import StateGraph, START, END, MessagesState
from langgraph.prebuilt import ToolNode, tools_condition
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.messages import SystemMessage, HumanMessage, RemoveMessage, AIMessage
from langchain_core.runnables import RunnableConfig
from langgraph.checkpoint.postgres.aio import AsyncPostgresSaver
from psycopg_pool import AsyncConnectionPool
from contextlib import asynccontextmanager
from typing_extensions import Literal, List, Annotated, Optional
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from dotenv import find_dotenv, load_dotenv
import os
import json
import io
import uuid
import logging
from src.backend.db import get_db, Session, Chats, User

# -----------------------------------------------------------------------------
# Structured Logging Setup
# -----------------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s"
)
logger = logging.getLogger("chatbot_backend")

# -----------------------------------------------------------------------------
# Configuration
# -----------------------------------------------------------------------------
load_dotenv(find_dotenv())

gemini_api_key = os.environ.get("GEMINI_API_KEY")
gemini_model_name = os.environ.get("GEMINI_MODEL", "gemini-3.5-flash-lite")
DB_URI = os.environ.get("DB_URI")


class ChatState(MessagesState):
    summary: str


# -----------------------------------------------------------------------------
# LangGraph Workflow Construction
# -----------------------------------------------------------------------------
def build_graph(checkpointer: AsyncPostgresSaver):
    def search_web(query: str) -> str:
        """Use this function to search the web for answers."""
        logger.info(f"Model calling search_web tool with query: '{query}'")
        return f"The search for '{query}' returned: Global software market projected to reach $687.9B in 2025."

    model = ChatGoogleGenerativeAI(
        model=gemini_model_name,
        api_key=gemini_api_key
    ).bind_tools([search_web])

    def call_model(state: ChatState, config: RunnableConfig) -> ChatState:
        summary = state.get("summary", "")
        if summary:
            sys_msg = f"Summary of conversation earlier: {summary}"
            messages = [SystemMessage(content=sys_msg)] + state["messages"]
        else:
            messages = state["messages"]
        response = model.invoke(messages)
        return {"messages": [response]}

    def summarize_conversation(state: ChatState, config: RunnableConfig):
        summary = state.get("summary", "")
        if summary:
            summary_message = (
                f"This is summary of the conversation to date: {summary}\n\n"
                "Extend the summary by taking into account the new messages above:"
            )
        else:
            summary_message = "Create a summary of the conversation above:"

        messages = state["messages"] + [HumanMessage(content=summary_message)]
        response = model.invoke(messages)
        delete_messages = [RemoveMessage(id=m.id) for m in state["messages"][:-2]]
        return {"summary": response.content, "messages": delete_messages}

    def should_continue(state: ChatState) -> Literal["summarize_conversation", "tools", "__end__"]:
        if tools_condition(state) == "tools":
            return "tools"
        if len(state["messages"]) > 20:
            return "summarize_conversation"
        return "__end__"

    workflow: StateGraph = StateGraph(ChatState)
    workflow.add_node("conversation", call_model)
    workflow.add_node("summarize_conversation", summarize_conversation)
    workflow.add_node("tools", ToolNode([search_web]))

    workflow.add_edge(START, "conversation")
    workflow.add_conditional_edges(
        "conversation",
        should_continue,
        {
            "summarize_conversation": "summarize_conversation",
            "tools": "tools",
            "__end__": END
        }
    )
    workflow.add_edge("summarize_conversation", END)
    workflow.add_edge("tools", "conversation")

    return workflow.compile(checkpointer)


# -----------------------------------------------------------------------------
# Modern FastAPI Lifespan Context Manager
# -----------------------------------------------------------------------------
@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info("Initializing connection pool and checkpointer...")
    connection_kwargs = {"autocommit": True, "prepare_threshold": 0}
    pool = AsyncConnectionPool(conninfo=DB_URI, max_size=20, kwargs=connection_kwargs)
    await pool.open()

    checkpointer = AsyncPostgresSaver(pool)
    await checkpointer.setup()

    graph = build_graph(checkpointer)

    # Attach state to FastAPI app
    app.state.pool = pool
    app.state.checkpointer = checkpointer
    app.state.graph = graph

    logger.info("Application setup complete. Ready to serve requests.")
    yield

    logger.info("Shutting down: closing database connection pool...")
    await pool.close()
    logger.info("Database connection pool closed.")


# -----------------------------------------------------------------------------
# App Initialization
# -----------------------------------------------------------------------------
app: FastAPI = FastAPI(
    title="Gemini Chatbot Backend Engine",
    description="Enterprise FastAPI backend for conversational streaming multimodal AI using Gemini & LangGraph.",
    lifespan=lifespan
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=[os.environ.get("ALLOWED_ORIGINS", "http://localhost:3000")],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)


# -----------------------------------------------------------------------------
# Document Processing Helper (100% In-Memory Parsing)
# -----------------------------------------------------------------------------
async def process_files(files: List[UploadFile] = File(...)):
    message_parts = []
    for file in files:
        content_type = file.content_type
        content = await file.read()

        if content_type.startswith("application/pdf"):
            reader = PdfReader(io.BytesIO(content))
            text = "Data from PDF : " + "\n".join(page.extract_text() for page in reader.pages)
            message_parts.append({"type": "text", "text": text})

        elif content_type.startswith("application/vnd.openxmlformats-officedocument.wordprocessingml.document"):
            doc = Document(io.BytesIO(content))
            text = "Data from Docx : " + "\n".join(p.text for p in doc.paragraphs)
            message_parts.append({"type": "text", "text": text})

        elif content_type.startswith("application/vnd.openxmlformats-officedocument.presentationml.presentation"):
            presentation = Presentation(io.BytesIO(content))
            text_runs = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
            text = "Data from PPTx : \n" + "\n".join(text_runs)
            message_parts.append({"type": "text", "text": text})
        else:
            text = "Data from File : \n" + content.decode("utf-8", errors="ignore")
            message_parts.append({"type": "text", "text": text})
    return message_parts


# -----------------------------------------------------------------------------
# Routes & Endpoints
# -----------------------------------------------------------------------------
@app.get("/")
async def read_root(request: Request):
    return {
        "status": "online",
        "service": "Gemini Multimodal AI Assistant Backend",
        "model": gemini_model_name,
        "graph_initialized": hasattr(request.app.state, "graph")
    }


@app.get("/chat/history/{thread_id}")
async def get_chat_history(request: Request, thread_id: str, db: Annotated[Session, Depends(get_db)]):
    user_id = request.cookies.get("user_id")
    logger.info(f"Fetching chat history for user_id={user_id}, thread_id={thread_id}")

    existing_chat = db.query(Chats).filter_by(thread_id=thread_id, user_id=user_id).first()
    if not existing_chat:
        return {"messages": []}

    checkpointer: AsyncPostgresSaver = request.app.state.checkpointer
    state = await checkpointer.aget({"configurable": {"thread_id": thread_id}})
    if not state:
        return {"messages": []}

    raw_messages = state["channel_values"]["messages"]
    simplified_messages = []
    for m in raw_messages:
        if isinstance(m, HumanMessage):
            content = m.content[0] if isinstance(m.content, list) else m.content
            simplified_messages.append({
                "role": "human",
                "content": content["text"] if isinstance(content, dict) else str(content)
            })
        elif isinstance(m, AIMessage):
            simplified_messages.append({"role": "ai", "content": str(m.content)})
        else:
            simplified_messages.append({"role": "system", "content": str(m.content)})

    return {"messages": simplified_messages}


@app.get("/all_chats")
async def get_all_chats(request: Request, db: Annotated[Session, Depends(get_db)]):
    user_id = request.cookies.get("user_id")
    if not user_id:
        raise HTTPException(status_code=401, detail="Not logged in")

    all_chats = db.query(Chats).filter_by(user_id=user_id).all()
    user = db.query(User).filter_by(id=user_id).first()
    user_email = user.email if user else "Guest"
    total_chats = len(all_chats)

    return {
        "chats": all_chats,
        "user_id": user_id,
        "user_email": user_email,
        "total_chats": total_chats
    }


@app.post("/chat_stream")
async def chat_stream(
    request: Request,
    query: str = Form(...),
    thread_id: str = Form(...),
    files: List[UploadFile] = File(None)
):
    config = {"configurable": {"thread_id": thread_id}}
    messages = [{"type": "text", "text": query}]
    if files:
        extracted_parts = await process_files(files)
        messages.extend(extracted_parts)

    graph = request.app.state.graph

    async def event_generator():
        try:
            node_to_stream = "conversation"
            input_message = HumanMessage(content=messages)
            async for event in graph.astream_events({"messages": [input_message]}, config, version="v2"):
                if event["event"] == "on_chat_model_stream" and event["metadata"].get("langgraph_node", "") == node_to_stream:
                    chunk = event["data"]["chunk"]
                    if chunk.content:
                        yield json.dumps({"type": "content", "response": chunk.content}) + "\n"
        except Exception as e:
            logger.error(f"Stream error: {e}")
            yield json.dumps({"type": "error", "message": str(e)}) + "\n"

    return StreamingResponse(event_generator(), media_type="application/x-ndjson")


@app.post("/new_chat_stream")
async def new_chat_stream(
    request: Request,
    db: Annotated[Session, Depends(get_db)],
    query: str = Form(...),
    files: Optional[List[UploadFile]] = File(None)
):
    user_id = request.cookies.get("user_id")
    messages = [{"type": "text", "text": query}]
    if files:
        extracted_parts = await process_files(files)
        messages.extend(extracted_parts)

    thread_id = str(uuid.uuid4())
    config = {"configurable": {"thread_id": thread_id}}
    graph = request.app.state.graph

    async def event_generator():
        try:
            yield json.dumps({"type": "init", "thread_id": thread_id}) + "\n"
            node_to_stream = "conversation"
            input_message = HumanMessage(content=messages)
            async for event in graph.astream_events({"messages": [input_message]}, config, version="v2"):
                if event["event"] == "on_chat_model_stream" and event["metadata"].get("langgraph_node", "") == node_to_stream:
                    chunk = event["data"]["chunk"]
                    if chunk and chunk.content:
                        yield json.dumps({"type": "content", "response": chunk.content}) + "\n"
                    else:
                        yield json.dumps({"type": "keepalive"}) + "\n"
        except Exception as e:
            logger.error(f"Error in event_generator: {e}")
            yield json.dumps({"type": "error", "message": str(e)}) + "\n"
            raise

        prompt = (
            f"Return ONLY one short, descriptive title (max 5 words) for this chat. "
            f"No bullet points, no quotes, no explanations. "
            f"First message: '{query}'"
        )
        temp_model = ChatGoogleGenerativeAI(model=gemini_model_name, api_key=gemini_api_key)
        temp_res = await temp_model.ainvoke(prompt)
        chat_name = temp_res.content.strip()

        yield json.dumps({"type": "final", "chat_name": chat_name, "thread_id": thread_id}) + "\n"

        temp_chat = Chats(thread_id=thread_id, user_id=user_id, chat_name=chat_name)
        db.add(temp_chat)
        db.commit()
        db.refresh(temp_chat)

    return StreamingResponse(event_generator(), media_type="application/x-ndjson")


@app.get("/rename_chat")
async def rename_chat(thread_id: str, new_name: str, db: Annotated[Session, Depends(get_db)]):
    chat = db.query(Chats).filter_by(thread_id=thread_id).first()
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    chat.chat_name = new_name
    db.commit()
    db.refresh(chat)
    return {"message": "Chat renamed successfully", "chat": chat}


@app.get("/delete_chat")
async def delete_chat(thread_id: str, db: Annotated[Session, Depends(get_db)]):
    chat = db.query(Chats).filter_by(thread_id=thread_id).first()
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    db.delete(chat)
    db.commit()
    return {"message": "Chat deleted successfully", "chat": chat}


# -----------------------------------------------------------------------------
# Auth Routes
# -----------------------------------------------------------------------------
@app.post("/auth/login")
def login(db: Annotated[Session, Depends(get_db)], email: str = Body(...), password: str = Body(...)):
    try:
        res = get_user(db=db, email=email, password=password)
        if "error" in res or "Error" in res or "Exception" in res:
            return JSONResponse(content={"error": str(res.get("error") or res.get("Exception")), "type": res.get("type", "exception")}, status_code=400)
        response_data = {"message": "Login successful"}
        response = JSONResponse(content=response_data)
        response.set_cookie(key="token", value=res["access_token"], httponly=True, secure=True, samesite="none")
        response.set_cookie(key="user_id", value=res["user"].id, httponly=True, secure=True, samesite="none")
        return response
    except Exception as e:
        return JSONResponse(content={"error": str(e), "type": "exception"}, status_code=500)


@app.post("/auth/signup")
def signup(db: Annotated[Session, Depends(get_db)], email: str = Body(...), password: str = Body(...), name: str | None = Body(None)):
    try:
        res = signup_user(db=db, email=email, password=password, name=name)
        if "error" in res or "Error" in res or "Exception" in res:
            return JSONResponse(content={"error": str(res.get("error") or res.get("Exception")), "type": res.get("type", "exception")}, status_code=400)
        response_data = {"message": "account created successfully"}
        response = JSONResponse(content=response_data)
        response.set_cookie(key="token", value=res["access_token"], httponly=True, secure=True, samesite="none")
        response.set_cookie(key="user_id", value=res["new_user"].id, httponly=True, secure=True, samesite="none")
        return response
    except Exception as e:
        return JSONResponse(content={"error": str(e), "type": "exception"}, status_code=500)


@app.get("/guest_login")
def guest_login(db: Annotated[Session, Depends(get_db)]):
    email = f"guest_{uuid.uuid4()}"
    password = "guest"
    try:
        res = signup_user(db=db, email=email, password=password)
        if "error" in res or "Error" in res or "Exception" in res:
            return JSONResponse(content={"error": str(res.get("error") or res.get("Exception")), "type": res.get("type", "exception")}, status_code=400)

        target_url = os.environ.get("ALLOWED_ORIGINS", "http://localhost:3000")
        response = RedirectResponse(target_url)
        response.set_cookie(key="token", value=res["access_token"], httponly=True, secure=True, samesite="none")
        response.set_cookie(key="user_id", value=res["new_user"].id, httponly=True, secure=True, samesite="none")
        return response
    except Exception as e:
        return JSONResponse(content={"error": str(e), "type": "exception"}, status_code=500)


@app.get("/auth/logout")
def logout():
    response = JSONResponse(content={"message": "Logged out successfully"})
    response.delete_cookie(key="token", path="/", secure=True, samesite="none")
    response.delete_cookie(key="user_id", path="/", secure=True, samesite="none")
    logger.info("User logged out, cookies cleared.")
    return response


# -----------------------------------------------------------------------------
# Local Dev Runner
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("src.backend.main:app", host="0.0.0.0", port=8000, reload=True)
